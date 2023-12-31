import sys
from PySide6.QtWidgets import QApplication, QMainWindow, QFileDialog, QPushButton, QVBoxLayout, QWidget, QMessageBox
from pptx import Presentation
from docx import Document
import PyPDF2
import chardet
import pandas as pd


class Editor(QMainWindow):
    def __init__(self):
        super().__init__()
        # self.excel_data = {}
        self.init_ui()

    def init_ui(self):
        self.central_widget = QWidget(self)
        self.setGeometry(100, 100, 200, 200)
        self.setWindowTitle('读写文件')

        self.setCentralWidget(self.central_widget)
        self.layout = QVBoxLayout(self.central_widget)

        # Open button
        open_button = QPushButton('打开文件', self)
        open_button.clicked.connect(self.open_file)
        self.layout.addWidget(open_button)

        # Save button
        save_button = QPushButton('保存文件', self)
        save_button.clicked.connect(self.save_file)
        self.layout.addWidget(save_button)

        # Store the opened file extension and content
        self.opened_file_extension = None
        self.content = ""

    def open_file(self):
        # Reset previous data

        self.reset_data()
        file_name, _ = QFileDialog.getOpenFileName(self, "Open File", "",
                                                   "EXECL Files (*.xlsx *.xls *.csv);;Text Files (*.txt);;PPTX Files (*.pptx);;DOCX Files (*.docx);;PDF Files (*.pdf)")

        if file_name:
            self.opened_file_extension = file_name.split('.')[-1].lower()

            if self.opened_file_extension == 'txt':
                with open(file_name, 'rb') as file:
                    result = chardet.detect(file.read())
                    encoding = result['encoding']
                with open(file_name, 'r', encoding=encoding) as file:
                    try:
                        self.content = file.read()
                    except UnicodeDecodeError:
                        # Handle the exception or try a different encoding
                        QMessageBox.warning(self, "Error", "Error decoding the file. Try a different encoding.")

            elif self.opened_file_extension == 'pptx':
                self.presentation = Presentation(file_name)
                for slide in self.presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            self.content += shape.text + "\n"

            elif self.opened_file_extension == 'docx':
                doc = Document(file_name)
                for paragraph in doc.paragraphs:
                    self.content += paragraph.text + "\n"

            elif self.opened_file_extension == 'pdf':
                self.pdf_file = open(file_name, 'rb')
                self.pdf_reader = PyPDF2.PdfReader(self.pdf_file)

            elif self.opened_file_extension == 'xlsx' or self.opened_file_extension == 'xls' or self.opened_file_extension == 'csv':
                try:

                    if self.opened_file_extension == 'csv':
                        self.excel_data = pd.read_csv(file_name, encoding="gbk")
                    else:
                        self.excel_data = pd.read_excel(file_name, sheet_name=None)
                    QMessageBox.information(self, "Success", "Excel file loaded successfully.")
                except Exception as e:
                    QMessageBox.warning(self, "Error", f"An error occurred during Excel file loading: {str(e)}")

    def save_file(self):
        if self.opened_file_extension is None:
            return
        try:
            if self.opened_file_extension == "txt":
                file_name, _ = QFileDialog.getSaveFileName(self, "Save File", "",
                                                           f"{self.opened_file_extension.upper()} Files (*.{self.opened_file_extension})")
                if file_name:
                    with open(file_name, 'w') as file:
                        file.write(self.content)
                    QMessageBox.information(self, "Success", "File saved successfully.")

            elif self.opened_file_extension == "pptx":
                output_ppt_path, _ = QFileDialog.getSaveFileName(self, "Save File", "",
                                                                 f"{self.opened_file_extension.upper()} Files (*.{self.opened_file_extension})")
                self.presentation.save(output_ppt_path)
                QMessageBox.information(self, "Success", "File saved successfully.")

            elif self.opened_file_extension == "docx":
                output_docx_path, _ = QFileDialog.getSaveFileName(self, "Save File", "",
                                                                  f"{self.opened_file_extension.upper()} Files (*.{self.opened_file_extension})")
                doc = Document()
                for line in self.content.split("\n"):
                    doc.add_paragraph(line)
                doc.save(output_docx_path)
                QMessageBox.information(self, "Success", "File saved successfully.")

            elif self.opened_file_extension == "pdf":
                output_pdf_path, _ = QFileDialog.getSaveFileName(self, "Save File", "",
                                                                 f"{self.opened_file_extension.upper()} Files (*.{self.opened_file_extension})")
                pdf_writer = PyPDF2.PdfWriter()

                # Add pages to the writer
                for page_num in range(len(self.pdf_reader.pages)):
                    page = self.pdf_reader.pages[page_num]
                    pdf_writer.add_page(page)

                # Write the new PDF file
                with open(output_pdf_path, 'wb') as output_file:
                    pdf_writer.write(output_file)
                self.pdf_file.close()
                QMessageBox.information(self, "Success", "File saved successfully.")

            elif self.opened_file_extension in ["xlsx", "xls", "csv"]:
                output_excel_path, _ = QFileDialog.getSaveFileName(self, "Save File", "",
                                                                   f"{self.opened_file_extension.upper()} Files (*.{self.opened_file_extension})")
                if output_excel_path.lower().endswith(".xlsx"):
                    with pd.ExcelWriter(output_excel_path,engine="xlrd",mode="a") as writer:
                        for sheet_name, df in self.excel_data.items():
                            df.to_excel(writer, sheet_name=sheet_name, index=False)
                    QMessageBox.information(self, "Success", "Excel file saved successfully.")
                elif output_excel_path.lower().endswith(".csv"):
                    self.excel_data.to_csv(output_excel_path, index=False)
                    QMessageBox.information(self, "Success", "Excel file saved successfully.")
                else:
                    with pd.ExcelWriter(output_excel_path,engine="openpyxl") as writer:
                        for sheet_name, df in self.excel_data.items():
                            df.to_excel(writer, sheet_name=sheet_name, index=False)


        except Exception as e:
            QMessageBox.warning(self, "Error", f"An error occurred during save: {str(e)}")

    def reset_data(self):
        # Reset data before opening a new file
        self.opened_file_extension = None
        self.content = ""


if __name__ == '__main__':
    app = QApplication(sys.argv)
    main_window = Editor()
    main_window.show()
    sys.exit(app.exec())
